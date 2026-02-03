import os
import re  # 用于提取文件名中的数字
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def extract_number(filename):
    """从文件名中提取数字（用于排序）"""
    # 正则匹配文件名中的数字（如从"Huawei_Blob_Img_8.png"中提取8）
    match = re.search(r'(\d+)', filename)
    if match:
        return int(match.group(1))  # 返回提取的数字
    return 0  # 无数字时返回0（排在最前）


def images_to_ppt(image_folder, output_ppt_path):
    """将文件夹中的所有图片按标号顺序批量插入PPT，每张PPT放一张图片"""
    # 1. 初始化PPT
    prs = Presentation()
    # 设置幻灯片大小为16:9（宽13.333英寸，高7.5英寸）
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 2. 获取文件夹中所有图片文件
    image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp')
    image_files = []

    for filename in os.listdir(image_folder):
        if filename.lower().endswith(image_extensions):
            image_path = os.path.join(image_folder, filename)
            image_files.append(image_path)

    # 关键修改：按文件名中的数字标号排序（而非字符串排序）
    # 排序依据：从文件名提取的数字（如1,2,3...10,11）
    image_files.sort(key=lambda x: extract_number(os.path.basename(x)))

    if not image_files:
        print("❌ 未找到任何图片文件！")
        return

    # 3. 遍历图片，插入幻灯片（此时已按标号顺序排列）
    for i, image_path in enumerate(image_files, 1):
        # 创建空白幻灯片
        slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题（显示图片文件名和实际标号）
        image_name = os.path.basename(image_path)
        title = slide.shapes.title
        title.text = f"图片 {i}/{len(image_files)}: {image_name}"
        title.text_frame.paragraphs[0].font.size = Pt(14)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 计算图片位置和大小
        margin = Inches(1)
        max_width = prs.slide_width - 2 * margin
        max_height = prs.slide_height - 2 * margin - Inches(1)

        # 插入图片
        left = margin
        top = margin + Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, width=None, height=None)

        # 按比例缩放（确保尺寸为整数）
        if pic.width > max_width:
            scale = max_width / pic.width
            pic.width = int(max_width)
            pic.height = int(pic.height * scale)
        if pic.height > max_height:
            scale = max_height / pic.height
            pic.height = int(max_height)
            pic.width = int(pic.width * scale)

        # 居中图片
        pic.left = (prs.slide_width - pic.width) // 2

        print(f"✅ 已插入第 {i} 张图片：{image_name}")

    # 4. 保存PPT
    prs.save(output_ppt_path)
    print(f"\n🎉 PPT生成完成！保存路径：{output_ppt_path}")


if __name__ == "__main__":
    # 图片文件夹路径（替换为你的路径）
    IMAGE_FOLDER = r"D:\PythonCode\PyCrawler\small_project\personalTest\Imgs"
    # 生成的PPT保存路径（替换为你的路径）
    OUTPUT_PPT = r"C:\Users\Yao\Desktop\图片汇总.pptx"

    images_to_ppt(IMAGE_FOLDER, OUTPUT_PPT)